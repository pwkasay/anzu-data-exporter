import io
import json
import logging
import os
import ssl
import time
from datetime import datetime, timezone
from io import BytesIO
import aiohttp
import asyncio

import PyPDF2
import certifi
import docx
import openai
import pandas as pd
import requests
from dateutil import parser
from dateutil.relativedelta import relativedelta
from pptx import Presentation

import uuid


def get_secrets():
    try:
        CLIENT_ID = os.getenv("CLIENT_ID")
        CLIENT_SECRET = os.getenv("CLIENT_SECRET")
        TENANT_ID = os.getenv("TENANT_ID")
        HUBSPOT_API_KEY = os.getenv("HUBSPOT_API_KEY")
        OPEN_AI_KEY = os.getenv("OPEN_AI_KEY")
        USER_ID = os.getenv("USER_ID")
        print(USER_ID)
        return (
            CLIENT_ID,
            CLIENT_SECRET,
            TENANT_ID,
            HUBSPOT_API_KEY,
            OPEN_AI_KEY,
            USER_ID,
        )
    except Exception as e:
        logging.error(f"Failed to retrieve secrets: {e}")
        raise


from dotenv import load_dotenv

# Determine environment mode from environment variable or default to 'dev'
mode = os.getenv("ENVIRONMENT", "dev")
if mode == "dev":
    load_dotenv()

(
    CLIENT_ID,
    CLIENT_SECRET,
    TENANT_ID,
    HUBSPOT_API_KEY,
    OPEN_AI_KEY,
    USER_ID,
) = get_secrets()
print("Env Setup")


def search_hubspot_object(object_type, search_body):
    url = f"https://api.hubapi.com/crm/v3/objects/{object_type}/search"
    headers = {
        "Content-Type": "application/json",
        "Authorization": f"Bearer {HUBSPOT_API_KEY}",
    }
    response = requests.post(url, headers=headers, data=json.dumps(search_body))
    response.raise_for_status()
    return response.json()


# Cache for storing owner details
owner_details_cache = {}


# Async function to fetch owner details with throttling and caching (v3 API)
async def fetch_owner_details_async(session, owner_id):
    if owner_id in owner_details_cache:
        return owner_id, owner_details_cache[owner_id]

    url = f"https://api.hubapi.com/crm/v3/owners/{owner_id}"
    headers = {
        "Content-Type": "application/json",
        "Authorization": f"Bearer {HUBSPOT_API_KEY}",
    }

    max_retries = 3
    backoff_time = 2

    for attempt in range(max_retries):
        try:
            async with session.get(url, headers=headers) as response:
                if response.status == 200:
                    owner_data = await response.json()
                    owner_details_cache[owner_id] = owner_data
                    return owner_id, owner_data
                elif response.status == 429:  # Rate limit
                    await asyncio.sleep(backoff_time)
                    backoff_time *= 2
                else:
                    logging.error(
                        f"Failed to fetch owner {owner_id}: {response.status}"
                    )
                    return owner_id, None
        except Exception as e:
            logging.error(f"Error fetching owner {owner_id}: {e}")
            if attempt == max_retries - 1:
                return owner_id, None
            await asyncio.sleep(backoff_time)

    return owner_id, None


# Sync wrapper for backward compatibility
def fetch_owner_details(owner_id):
    if owner_id in owner_details_cache:
        return owner_details_cache[owner_id]
    url = f"https://api.hubapi.com/crm/v3/owners/{owner_id}"
    headers = {
        "Content-Type": "application/json",
        "Authorization": f"Bearer {HUBSPOT_API_KEY}",
    }
    retries = 3
    while retries > 0:
        try:
            response = requests.get(url, headers=headers)
            response.raise_for_status()
            owner_details_cache[owner_id] = response.json()
            return owner_details_cache[owner_id]
        except requests.exceptions.RequestException as e:
            logging.error(f"Failed to fetch owner details for ID {owner_id}: {e}")
            retries -= 1
            if retries == 0:
                raise
            time.sleep(2 ** (3 - retries))  # Exponential backoff


# Async function to fetch and attach owner details
async def fetch_and_attach_owner_details_async(deals, owner_property, batch_size=10):
    owner_ids = [
        deal["properties"].get(owner_property)
        for deal in deals
        if deal["properties"].get(owner_property)
    ]
    unique_owner_ids = list(set(owner_ids))

    sslcontext = ssl.create_default_context(cafile=certifi.where())
    async with aiohttp.ClientSession(
        connector=aiohttp.TCPConnector(ssl=sslcontext)
    ) as session:
        owners = {}

        # Process in batches to avoid overwhelming the API
        for i in range(0, len(unique_owner_ids), batch_size):
            batch = unique_owner_ids[i : i + batch_size]
            tasks = [fetch_owner_details_async(session, owner_id) for owner_id in batch]
            results = await asyncio.gather(*tasks)

            for owner_id, owner_data in results:
                if owner_data:
                    owners[owner_id] = owner_data

            # Small delay between batches
            if i + batch_size < len(unique_owner_ids):
                await asyncio.sleep(0.5)

    # Attach owner details to deals
    for deal in deals:
        owner_id = deal["properties"].get(owner_property)
        if owner_id:
            deal[f"{owner_property}_details"] = owners.get(owner_id, {})

    return deals


# Sync wrapper for backward compatibility
def fetch_and_attach_owner_details(deals, owner_property):
    return asyncio.run(fetch_and_attach_owner_details_async(deals, owner_property))


def download_file(url):
    response = requests.get(url)
    response.raise_for_status()
    return response.content


def read_pdf(file_content):
    pdf_reader = PyPDF2.PdfReader(BytesIO(file_content))
    text = ""
    for page_num in range(len(pdf_reader.pages)):
        page = pdf_reader.pages[page_num]
        text += page.extract_text()
    return text


def read_word(file_content):
    doc = docx.Document(BytesIO(file_content))
    text = "\n".join([paragraph.text for paragraph in doc.paragraphs])
    return text


def read_excel(file_content):
    excel_file = pd.ExcelFile(BytesIO(file_content))
    sheets = {}
    for sheet_name in excel_file.sheet_names:
        sheets[sheet_name] = excel_file.parse(sheet_name)
    return sheets


def read_ppt(file_content):
    presentation = Presentation(BytesIO(file_content))
    text = ""
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text


def update_file_access(file_id, access_level):
    url = f"https://api.hubapi.com/files/v3/files/{file_id}"
    headers = {
        "Authorization": f"Bearer {HUBSPOT_API_KEY}",
        "Content-Type": "application/json",
    }
    data = {"options": {"access": access_level}}
    response = requests.patch(url, headers=headers, data=json.dumps(data))
    response.raise_for_status()
    return response.json()


def generate_signed_url(file_id):
    url = f"https://api.hubapi.com/files/v3/files/{file_id}/signed-url"
    headers = {
        "Authorization": f"Bearer {HUBSPOT_API_KEY}",
        "Content-Type": "application/json",
    }
    response = requests.get(url, headers=headers)
    response.raise_for_status()
    return response.json()["url"]


def get_file_details(file_id):
    url = f"https://api.hubapi.com/files/v3/files/{file_id}"
    headers = {"Authorization": f"Bearer {HUBSPOT_API_KEY}"}
    response = requests.get(url, headers=headers)
    response.raise_for_status()
    return response.json()


failed_attachments = []


def fetch_deal_properties():
    """Fetch deal properties using v3 API."""
    properties_url = "https://api.hubapi.com/crm/v3/properties/deals"
    headers = {
        "Authorization": f"Bearer {HUBSPOT_API_KEY}",
        "Content-Type": "application/json",
    }
    response = requests.get(properties_url, headers=headers)
    response.raise_for_status()
    return response.json().get("results", [])


def fetch_single_deal_with_history(deal_id):
    properties_url = "https://api.hubapi.com/properties/v2/deals/properties"
    headers = {
        "Authorization": f"Bearer {HUBSPOT_API_KEY}",
        "Content-Type": "application/json",
    }
    # Fetch all deal properties
    response = requests.get(properties_url, headers=headers)
    if response.status_code == 200:
        deal_properties = response.json()
        all_properties = [prop["name"] for prop in deal_properties]
        # HubSpot API endpoint for fetching a single deal with property history
        deal_url = f"https://api.hubapi.com/crm/v3/objects/deals/{deal_id}"
        params = {
            "propertiesWithHistory": ",".join(
                all_properties
            )  # Fetch all properties with their history
        }
        deal_response = requests.get(deal_url, headers=headers, params=params)
        if deal_response.status_code == 200:
            deal_info = deal_response.json()
            print(deal_info)
        else:
            print(f"Error: {deal_response.status_code}")
            print(deal_response.text)
    else:
        print(f"Error: {response.status_code}")
        print(response.text)


def fetch_single_deal(deal_id):
    properties_url = "https://api.hubapi.com/properties/v2/deals/properties"
    url = f"https://api.hubapi.com/crm/v3/objects/deals/{deal_id}"
    headers = {
        "Authorization": f"Bearer {HUBSPOT_API_KEY}",
        "Content-Type": "application/json",
    }
    params = {
        "properties": [
            "dealname",
            "amount",
            "dealstage",
            "pipeline",
            "closedate",
            "createdate",
            "hubspot_owner_id",
            "fund",
        ]
    }
    properties_response = requests.get(properties_url, headers=headers)
    if properties_response.status_code == 200:
        deal_properties = properties_response.json()
        all_properties = [prop["name"] for prop in deal_properties]
    else:
        print(f"Error: {properties_response.status_code}")
    # Make the GET request to fetch the deal information
    response = requests.get(url, headers=headers, params=params)
    # Check if the request was successful
    if response.status_code == 200:
        deal_info = response.json()
    else:
        print(f"Error: {response.status_code}")


# DEPRECATED - Use fetch_deals() with include_stage_history=True instead
def fetch_deals_with_stage_history(start_date=None, end_date=None):
    return fetch_deals(
        start_date=start_date, end_date=end_date, include_stage_history=True
    )


async def fetch_stage_history(session, deal_id, stage_mapping):
    url = f"https://api.hubapi.com/crm/v3/objects/deals/{deal_id}"
    headers = {
        "Authorization": f"Bearer {HUBSPOT_API_KEY}",
        "Content-Type": "application/json",
    }
    params = {"propertiesWithHistory": "dealstage"}
    max_retries = 5
    backoff_time = 2  # Initial backoff time in seconds
    for attempt in range(max_retries):
        async with session.get(url, headers=headers, params=params) as response:
            if response.status == 200:
                deal_stage_history = await response.json()
                deal_stage_history = deal_stage_history["propertiesWithHistory"][
                    "dealstage"
                ]
                for stage in deal_stage_history:
                    stage_id = stage["value"]
                    stage["stage_name"] = stage_mapping.get(stage_id, "Unknown Stage")
                return deal_id, deal_stage_history
            elif response.status == 429:  # Rate limit error
                print(
                    f"Rate limit hit: waiting {backoff_time} seconds before retrying..."
                )
                await asyncio.sleep(backoff_time)
                backoff_time *= 2  # Exponentially increase backoff time
            else:
                print(
                    f"Failed to retrieve deal data: {response.status} {await response.text()}"
                )
                return deal_id, None
    print(f"Max retries reached for deal {deal_id}.")
    return deal_id, None


async def fetch_all_stage_histories(
    deals, stage_mapping, batch_size=140, delay_between_batches=10
):
    sslcontext = ssl.create_default_context(cafile=certifi.where())
    async with aiohttp.ClientSession(
        connector=aiohttp.TCPConnector(ssl=sslcontext)
    ) as session:
        stage_histories = {}
        for i in range(0, len(deals), batch_size):
            batch = deals[i : i + batch_size]
            tasks = [
                fetch_stage_history(session, deal["id"], stage_mapping)
                for deal in batch
            ]
            results = await asyncio.gather(*tasks)
            stage_histories.update(
                {deal_id: history for deal_id, history in results if history}
            )
            # Delay between batches to avoid hitting rate limits
            await asyncio.sleep(delay_between_batches)
        return stage_histories


async def fetch_and_attach_notes(deals, batch_size=4, delay_between_batches=1):
    sslcontext = ssl.create_default_context(cafile=certifi.where())
    async with aiohttp.ClientSession(
        connector=aiohttp.TCPConnector(ssl=sslcontext)
    ) as session:
        notes = {}
        for i in range(0, len(deals), batch_size):
            batch = deals[i : i + batch_size]
            tasks = [fetch_notes(session, deal["id"]) for deal in batch]
            results = await asyncio.gather(*tasks)
            notes.update({deal_id: note for deal_id, note in results if note})
            # Delay between batches
            await asyncio.sleep(delay_between_batches)
        for deal in deals:
            deal_id = deal["id"]
            deal_notes = notes.get(deal_id)
            deal["notes"] = deal_notes["results"]
        return deals


# DEPRECATED - Use fetch_deals() instead
def fetch_all_deals(start_date=None, end_date=None):
    return fetch_deals(start_date=start_date, end_date=end_date)


def fetch_deals_and_stage_histories(start_date=None, end_date=None):
    """Fetch deals with stage histories and owner details - optimized version."""
    deals = fetch_deals(start_date, end_date)

    # Run both owner fetches in parallel using async
    async def fetch_all_owners():
        await asyncio.gather(
            fetch_and_attach_owner_details_async(deals, "hubspot_owner_id"),
            fetch_and_attach_owner_details_async(deals, "team_member_1"),
        )

    asyncio.run(fetch_all_owners())

    # Fetch all stage histories asynchronously with batching and delays
    stage_mapping = get_stage_mapping()
    stage_histories = asyncio.run(fetch_all_stage_histories(deals, stage_mapping))

    # Merge the stage histories into the deals
    for deal in deals:
        deal_id = deal["id"]
        deal["deal_stage_history"] = stage_histories.get(deal_id, [])

    return deals


def fetch_deals(start_date=None, end_date=None, include_stage_history=False):
    properties = [
        "dealname",
        "priority",
        "referral_type",
        "pipeline",
        "broad_category_updated",
        "subcategory",
        "fund",
        "hubspot_owner_id",
        "team_member_1",
        "createdate",
        "keywords",
    ]
    deals = []
    if start_date is None and end_date is None:
        start_date = str(datetime.now().date() + relativedelta(months=-3))
        end_date = str(datetime.now().date() + relativedelta(days=+1))
    start_date = datetime.strptime(start_date, "%Y-%m-%d")
    end_date = datetime.strptime(end_date, "%Y-%m-%d")
    # Hubspot needs datetime to be set to midnight
    start_date = start_date.replace(hour=0, minute=0, second=0, microsecond=0)
    end_date = end_date.replace(hour=0, minute=0, second=0, microsecond=0)
    # Convert the close date to a Unix timestamp in milliseconds
    start_date = int(start_date.timestamp() * 1000)
    end_date = int(end_date.timestamp() * 1000)
    after = None

    max_retries = 3
    base_delay = 0.5

    search_results = None

    while True:
        search_body = {
            "filterGroups": [
                {
                    "filters": [
                        {
                            "propertyName": "createdate",
                            "operator": "BETWEEN",
                            "highValue": end_date,
                            "value": start_date,
                        },
                        {
                            "propertyName": "pipeline",
                            "operator": "EQ",
                            "value": "default",
                        },
                    ]
                }
            ],
            "properties": properties,
            "limit": 100,
        }
        if after:
            search_body["after"] = after

        for attempt in range(max_retries):
            try:
                search_results = search_hubspot_object("deals", search_body)
                break
            except Exception as e:
                if attempt < max_retries - 1:
                    delay = base_delay * (attempt**2)
                    time.sleep(delay)
                else:
                    print(f"Attempt {attempt + 1} failed: {e}. No more retries left.")
                    raise

        fetched_deals = search_results.get("results", [])

        # If including stage history, fetch it inline
        if include_stage_history:
            stage_mapping = get_stage_mapping()
            for deal in fetched_deals:
                deal_id = deal["id"]
                url = f"https://api.hubapi.com/crm/v3/objects/deals/{deal_id}"
                headers = {
                    "Authorization": f"Bearer {HUBSPOT_API_KEY}",
                    "Content-Type": "application/json",
                }
                params = {"propertiesWithHistory": "dealstage"}
                response = requests.get(url, headers=headers, params=params)
                if response.status_code == 200:
                    deal_stage_history = response.json()["propertiesWithHistory"][
                        "dealstage"
                    ]
                    for stage in deal_stage_history:
                        stage_id = stage["value"]
                        stage["stage_name"] = stage_mapping.get(
                            stage_id, "Unknown Stage"
                        )
                    deal["deal_stage_history"] = deal_stage_history
                else:
                    logging.error(
                        f"Failed to retrieve stage history: {response.status_code}"
                    )
                    deal["deal_stage_history"] = []

        deals.extend(fetched_deals)
        pagination = search_results.get("paging", [])
        if pagination and "next" in pagination:
            after = pagination["next"]["after"]
        else:
            break
    return deals


def organize_deals(deals):
    deals_dict = {}
    for deal in deals:
        broad_category = deal["properties"]["broad_category_updated"]
        if broad_category not in deals_dict:
            deals_dict[broad_category] = deal
    return deals_dict


def organize_cleaned_deals(cleaned_deals):
    deals_dict = {}
    for deal in cleaned_deals:
        broad_category = deal["broad_category_updated"]
        if broad_category not in deals_dict:
            deals_dict[broad_category] = [deal]
        else:
            deals_dict[broad_category].append(deal)
    return deals_dict


def read_prompt_text(text_path):
    try:
        with open(text_path, "r") as file:
            return file.read()
    except FileNotFoundError:
        return "Flag as no prompt"


gpt_prompt_path = "data/gpt_prompt.txt"
gpt_prompt = read_prompt_text(gpt_prompt_path)


def create_openai_client(OPEN_AI_KEY):
    openai_client = openai.OpenAI(api_key=OPEN_AI_KEY)
    return openai_client


openai_client = create_openai_client(OPEN_AI_KEY)

gpt_errors = []


def parse_with_chatgpt(openai_client, deal):
    try:
        messages = [
            {"role": "system", "content": gpt_prompt},
            {
                "role": "user",
                "content": f"Deal info: {deal}",
            },
        ]
        response = openai_client.chat.completions.create(
            model="gpt-5",
            messages=messages,
            max_tokens=2500,
            n=1,
            stop=None,
            temperature=0.5,
        )
        return response.choices[0].message.content
    except Exception as e:
        gpt_errors.append({deal["id"]: str(e)})
        return None


user_prompt_for_final_path = "data/user_prompt_for_final.txt"
user_prompt_for_final = read_prompt_text(user_prompt_for_final_path)

system_prompt_for_final_path = "data/system_prompt_for_final.txt"
system_prompt_for_final = read_prompt_text(system_prompt_for_final_path)


def batch_with_chatgpt(openai_client, deals):
    jsonl_lines = []
    for deal in deals:
        prompt = {
            "custom_id": str(uuid.uuid4()),
            "method": "POST",
            "url": "/v1/chat/completions",
            "body": {
                "model": "gpt-4o",
                "response_format": {"type": "json_object"},
                "messages": [
                    {"role": "system", "content": gpt_prompt},
                    {"role": "user", "content": f"Deal info: {deal}"},
                ],
                "max_tokens": 2500,
                "n": 1,
                "stop": None,
                "temperature": 0.5,
            },
        }
        jsonl_lines.append(json.dumps(prompt).encode("utf-8"))
    json_memory_file = io.BytesIO()
    for line in jsonl_lines:
        json_memory_file.write(line + b"\n")
    json_memory_file.seek(0)
    batch_input_file = openai_client.files.create(
        file=json_memory_file, purpose="batch"
    )
    batch_input_file_id = batch_input_file.id
    batch = openai_client.batches.create(
        input_file_id=batch_input_file_id,
        endpoint="/v1/chat/completions",
        completion_window="24h",
        metadata={"description": "deal data recommendation generator"},
    )
    return batch


def check_gpt(openai_client, batch):
    # check status
    file_response = None
    retrieved_batch = openai_client.batches.retrieve(batch.id)
    if retrieved_batch.status == "completed" and retrieved_batch.output_file_id:
        file_response = openai_client.files.content(retrieved_batch.output_file_id)
        return file_response.content
    elif retrieved_batch.status == "completed" and retrieved_batch.error_file_id:
        file_response = openai_client.files.content(retrieved_batch.error_file_id)
        raise Exception(
            f"Batch processing failed. Error details: {file_response.content}"
        )
    elif retrieved_batch.status == "failed":
        raise Exception(f"Batch processing failed. Error details: {retrieved_batch}")
    else:
        return None


def poll_gpt_check(check):
    if isinstance(check, bytes):
        result = check
        memory_file = io.BytesIO()
        memory_file.write(result)
        memory_file.seek(0)
        memory_text = io.StringIO(memory_file.getvalue().decode("utf-8"))
        results = []
        # Read from the in-memory text stream
        for line in memory_text:
            json_object = json.loads(line.strip())
            results.append(json_object)
        return results


def delete_batch_file(openai_client, batch):
    # finish delete file

    all_batches = openai_client.batches.list(limit=10)
    # delete function
    files = openai_client.files.list()

    openai_client.files.delete(batch.output_file_id)


def compile_with_chatgpt(openai_client, cleaned_deals):
    try:
        messages = [
            {"role": "system", "content": system_prompt_for_final},
            {
                "role": "user",
                "content": f"Context: {user_prompt_for_final} - Input Data: {cleaned_deals}",
            },
        ]
        response = openai_client.chat.completions.create(
            model="gpt-5",
            messages=messages,
            max_tokens=2500,
            n=1,
            stop=None,
            temperature=0.5,
        )
        return response.choices[0].message.content
    except Exception as e:
        print(f"Error: {e}")
        return None


def add_field_to_deal(field):
    url = "https://api.hubapi.com/properties/v1/deals/properties"

    headers = {
        "Content-Type": "application/json",
        "Authorization": f"Bearer {HUBSPOT_API_KEY}",
    }

    label = field
    name = field.replace(" ", "_").lower()

    payload = {
        "name": "broad_category",
        "label": "Broad Category",
        "description": "The broad category of the deal",
        "groupName": "dealinformation",
        "type": "string",
        "fieldType": "text",
    }

    response = requests.post(url, json=payload, headers=headers)

    if response.status_code == 201:
        print("Custom property created successfully!")
    else:
        print(f"Failed to create property: {response.status_code}, {response.text}")


def update_deal(deal):
    update_url = f"https://api.hubapi.com/deals/v1/deal/{deal.id}"
    headers = {
        "Authorization": f"Bearer {HUBSPOT_API_KEY}",
        "Content-Type": "application/json",
    }

    update_payload = {
        "properties": [
            {"name": "broad_category_updated", "value": "value"},
            {"name": "subcategory", "value": "value"},
        ]
    }

    response = requests.put(update_url, json=update_payload, headers=headers)

    if response.status_code == 200:
        print("Deal updated successfully!")
    else:
        print(f"Failed to update deal: {response.status_code}, {response.text}")


def create_hubspot_field():
    pass


def update_hubspot_keywords(deal):
    update_url = f"https://api.hubapi.com/deals/v1/deal/{deal['id']}"
    headers = {
        "Authorization": f"Bearer {HUBSPOT_API_KEY}",
        "Content-Type": "application/json",
    }
    if deal.get("parsed"):
        if deal["parsed"].get("recommendation"):
            rec_keys = str(deal["parsed"]["recommendation"])
            update_payload = {"properties": [{"name": "keywords", "value": rec_keys}]}
            response = requests.put(update_url, json=update_payload, headers=headers)
            return response


def get_stage_mapping():
    headers = {
        "Authorization": f"Bearer {HUBSPOT_API_KEY}",
        "Content-Type": "application/json",
    }
    pipeline_url = f"https://api.hubapi.com/crm-pipelines/v1/pipelines/deals/default"
    pipeline_response = requests.get(pipeline_url, headers=headers)
    pipeline_data = pipeline_response.json()
    # Create a mapping of stage IDs to stage labels
    stage_mapping = {
        stage["stageId"]: stage["label"] for stage in pipeline_data["stages"]
    }
    # Store the mapping in the dictionary with the deal ID as the key
    return stage_mapping


def get_deal_stage_name(stage_id, pipeline_id):
    # Fetch the pipeline data using the pipeline ID
    headers = {
        "Authorization": f"Bearer {HUBSPOT_API_KEY}",
        "Content-Type": "application/json",
    }
    pipeline_url = (
        f"https://api.hubapi.com/crm-pipelines/v1/pipelines/deals/{pipeline_id}"
    )
    pipeline_response = requests.get(pipeline_url, headers=headers)
    if pipeline_response.status_code != 200:
        raise Exception(
            f"Failed to fetch pipeline data: {pipeline_response.status_code} {pipeline_response.text}"
        )
    pipeline_data = pipeline_response.json()
    # Search for the stage with the matching ID
    for stage in pipeline_data["stages"]:
        if stage["stageId"] == stage_id:
            return stage["label"]


def get_deal_stage_history(deals):
    stage_mapping = get_stage_mapping()
    for deal in deals:
        time.sleep(0.1)
        deal_id = deal["id"]
        # Define the API endpoint and parameters
        url = f"https://api.hubapi.com/crm/v3/objects/deals/{deal_id}"
        headers = {
            "Authorization": f"Bearer {HUBSPOT_API_KEY}",
            "Content-Type": "application/json",
        }
        params = {
            "propertiesWithHistory": "dealstage",
        }
        # Make the request to HubSpot API
        response = requests.get(url, headers=headers, params=params)
        if response.status_code == 200:
            deal_stage_history = response.json()["propertiesWithHistory"]["dealstage"]
            for stage in deal_stage_history:
                stage_id = stage["value"]
                if stage_mapping.get(stage_id):
                    stage_name = stage_mapping[stage_id]
                    stage["stage_name"] = stage_name
                else:
                    stage.pop(stage_id, None)
            deal["deal_stage_history"] = deal_stage_history
        else:
            print(
                f"Failed to retrieve deal data: {response.status_code} {response.text}"
            )
    return deals


async def fetch_notes_attachments_and_engagements(
    deals, batch_size=4, delay_between_batches=1.2
):
    sslcontext = ssl.create_default_context(cafile=certifi.where())

    async with aiohttp.ClientSession(
        connector=aiohttp.TCPConnector(ssl=sslcontext)
    ) as session:
        for i in range(0, len(deals), batch_size):
            batch = deals[i : i + batch_size]
            tasks = [fetch_data_for_deal(session, deal["id"]) for deal in batch]
            results = await asyncio.gather(*tasks)

            for deal, result in zip(batch, results):
                deal["notes"] = result.get("notes")
                deal["attachments"] = result.get("attachments")
                deal["engagements"] = result.get("engagements")

            await asyncio.sleep(delay_between_batches)
        return deals


async def fetch_data_for_deal(session, deal_id):
    notes_task = fetch_notes(session, deal_id)
    engagements_task = fetch_engagements(session, deal_id)

    notes_result, engagements_result = await asyncio.gather(
        notes_task, engagements_task
    )

    notes, attachments = await process_notes_for_attachments(
        session, notes_result["results"]
    )

    return {
        "notes": notes,
        "attachments": attachments,
        "engagements": engagements_result,
    }


async def fetch_notes(session, deal_id):
    url = "https://api.hubapi.com/crm/v3/objects/notes/search"
    headers = {
        "Content-Type": "application/json",
        "Authorization": f"Bearer {HUBSPOT_API_KEY}",
    }
    search_body = {
        "filterGroups": [
            {
                "filters": [
                    {
                        "propertyName": "associations.deal",
                        "operator": "EQ",
                        "value": deal_id,
                    }
                ]
            }
        ],
        "properties": ["hs_note_body", "hs_attachment_ids", "propertyName"],
    }

    max_retries = 5
    backoff_time = 2
    for attempt in range(max_retries):
        async with session.post(url, headers=headers, json=search_body) as response:
            if response.status == 200:
                return await response.json()
            elif response.status == 429:  # Rate limit error
                print(
                    f"Rate limit hit: waiting {backoff_time} seconds before retrying..."
                )
                await asyncio.sleep(backoff_time)
                backoff_time *= 2  # Exponentially increase backoff time
            else:
                print(
                    f"Failed to retrieve deal data: {response.status} {await response.text()}"
                )
                return None

    print(f"Max retries reached for deal {deal_id}.")
    return None


async def process_notes_for_attachments(session, notes):
    attachments = []
    tasks = []
    note_indices = []

    for note_idx, note in enumerate(notes):
        if note["properties"].get("hs_attachment_ids"):
            file_id = note["properties"]["hs_attachment_ids"]
            tasks.append(fetch_attachment(session, file_id))
            note_indices.append(note_idx)

    if tasks:
        results = await asyncio.gather(*tasks, return_exceptions=True)

        for note_idx, result in zip(note_indices, results):
            if isinstance(result, Exception):
                print(f"Failed to fetch attachment for note index {note_idx}: {result}")
            else:
                attachments.append(result)

    return notes, attachments


async def fetch_attachment(session, file_id):
    url = f"https://api.hubapi.com/filemanager/api/v3/files/{file_id}/signed-url"
    headers = {
        "Authorization": f"Bearer {HUBSPOT_API_KEY}",
        "Content-Type": "application/json",
        "User-Agent": "Custom",
    }
    async with session.get(url, headers=headers) as response:
        response.raise_for_status()
        data = await response.json()
        file_url = data["url"]
        extension = data["extension"]
        name = data["name"]

        async with session.get(file_url) as file_response:
            file_content = await file_response.read()

        if extension == "pdf":
            return read_pdf(file_content)
        elif extension == "docx":
            return read_word(file_content)
        elif extension == "xlsx":
            return read_excel(file_content)
        elif extension == "pptx":
            return read_ppt(file_content)
        else:
            raise ValueError(f"Unsupported file type: {extension} for {name}")


async def fetch_engagements(session, deal_id, engagement_type="EMAIL"):
    url = f"https://api.hubapi.com/engagements/v1/engagements/associated/deal/{deal_id}/paged"
    headers = {
        "Authorization": f"Bearer {HUBSPOT_API_KEY}",
        "Content-Type": "application/json",
    }
    params = {"limit": 100, "offset": 0}
    engagements = []

    while True:
        async with session.get(url, headers=headers, params=params) as response:
            response.raise_for_status()
            data = await response.json()
            engagements.extend(
                [
                    engagement
                    for engagement in data["results"]
                    if engagement["engagement"]["type"] == engagement_type
                ]
            )
            if not data.get("hasMore"):
                break
            params["offset"] = data["offset"]

    return engagements


def export_csv(start_date=None, end_date=None):
    """Export deals data to CSV with improved error handling and performance."""
    logging.info(f"Starting CSV export for dates: {start_date} to {end_date}")

    # Fetch deals with all enriched data
    deals = fetch_deals_and_stage_histories(start_date, end_date)
    logging.info(f"Fetched {len(deals)} deals")

    # Fetch and process fund mapping
    fund_mapping = {}
    try:
        fetched_deal_properties = fetch_deal_properties()
        for prop in fetched_deal_properties:
            if prop.get("name") == "fund" and "options" in prop:
                for option in prop["options"]:
                    fund_mapping[option.get("value")] = option.get(
                        "label", option.get("value")
                    )
    except Exception as e:
        logging.warning(f"Could not fetch fund mapping: {e}")

    flattened_data = []
    for deal in deals:
        # Start with properties but avoid full copy for memory efficiency
        flattened_entry = {}
        flattened_entry.update(deal.get("properties", {}))
        flattened_entry["id"] = deal.get("id")
        flattened_entry["createdAt"] = deal.get("createdAt")
        flattened_entry["updatedAt"] = deal.get("updatedAt")
        flattened_entry["archived"] = deal.get("archived", False)

        # Safe fund mapping with fallback
        if "fund" in flattened_entry and flattened_entry["fund"]:
            flattened_entry["fund"] = fund_mapping.get(
                flattened_entry["fund"],
                flattened_entry["fund"],  # Keep original if not in mapping
            )
        # Safe extraction of Lead Owner details
        if deal.get("hubspot_owner_id_details"):
            lead_owner = deal.get("hubspot_owner_id_details")
            first_name = lead_owner.get("firstName", "")
            last_name = lead_owner.get("lastName", "")
            flattened_entry["Lead Owner Name"] = f"{first_name} {last_name}".strip()
            flattened_entry["Lead Owner Email"] = lead_owner.get("email", "")
        else:
            flattened_entry["Lead Owner Name"] = ""
            flattened_entry["Lead Owner Email"] = ""

        # Safe extraction of Support Member details
        if deal.get("team_member_1_details"):
            support_member = deal.get("team_member_1_details")
            first_name = support_member.get("firstName", "")
            last_name = support_member.get("lastName", "")
            flattened_entry["Support Member Name"] = f"{first_name} {last_name}".strip()
            flattened_entry["Support Member Email"] = support_member.get("email", "")
        else:
            flattened_entry["Support Member Name"] = ""
            flattened_entry["Support Member Email"] = ""
        # Calculate time spent in each stage
        stage_durations = {}
        if "deal_stage_history" in deal and deal["deal_stage_history"]:
            stages = deal["deal_stage_history"]
            stages.sort(key=lambda x: x.get("timestamp", 0))  # Sort by timestamp

            # Process transitions between stages
            for i in range(len(stages) - 1):
                stage_name = stages[i].get("stage_name", "Unknown Stage")
                entry_timestamp = stages[i].get("timestamp")
                exit_timestamp = stages[i + 1].get("timestamp")

                try:
                    entry_time = parser.parse(entry_timestamp)
                    exit_time = parser.parse(exit_timestamp)
                    duration = (exit_time - entry_time).days
                    stage_durations[stage_name] = (
                        stage_durations.get(stage_name, 0) + duration
                    )
                except Exception as e:
                    logging.warning(
                        f"Error parsing timestamps for deal {deal.get('id')}: {e}"
                    )

            # Handle the last stage (current stage) - calculated ONCE outside the loop
            if stages:
                last_stage = stages[-1]
                last_stage_name = last_stage.get("stage_name", "Unknown Stage")
                last_entry_timestamp = last_stage.get("timestamp")

                try:
                    last_entry_time = parser.parse(last_entry_timestamp)
                    last_exit_time = datetime.now(timezone.utc)
                    last_duration = (last_exit_time - last_entry_time).days
                    stage_durations[last_stage_name] = (
                        stage_durations.get(last_stage_name, 0) + last_duration
                    )
                except Exception as e:
                    logging.warning(
                        f"Error parsing last stage timestamp for deal {deal.get('id')}: {e}"
                    )
        for stage_name, duration in stage_durations.items():
            flattened_entry[f"{stage_name}_days_in_stage"] = duration
        flattened_data.append(flattened_entry)
        # Convert to DataFrame
    df = pd.DataFrame(flattened_data)
    # List of columns to exclude (only drop if they exist)
    columns_to_exclude = [
        "hs_object_id",
        "archived",
        "hs_lastmodifieddate",
        "hubspot_owner_id",
        "pipeline",
        "team_member_1",
        "id",
        "createdAt",
        "updatedAt",
        "Lead Owner Email",
        "Support Member Email",
    ]
    # Drop only columns that exist in the dataframe
    columns_to_drop = [col for col in columns_to_exclude if col in df.columns]
    if columns_to_drop:
        df = df.drop(columns=columns_to_drop)
    if start_date is None and end_date is None:
        start_date = str(datetime.now().date() + relativedelta(months=-3))
        end_date = str(datetime.now().date() + relativedelta(days=+1))
    start_date = str(datetime.strptime(start_date, "%Y-%m-%d").date())
    end_date = str(datetime.strptime(end_date, "%Y-%m-%d").date())
    # Generate filename with proper date handling
    filename = f"Deal_Export--{start_date}-{end_date}.csv"

    # Create CSV in memory
    file_object = io.StringIO()
    df.to_csv(file_object, index=False)
    file_object.seek(0)

    # Log export metrics
    logging.info(f"CSV export completed: {len(df)} rows, {len(df.columns)} columns")
    logging.info(f"Filename: {filename}")

    return file_object, filename
