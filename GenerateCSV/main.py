import io
import json
import logging
import os
import time
from datetime import datetime
from io import BytesIO

import PyPDF2
import docx
import openai
import pandas as pd
import requests
from dateutil.relativedelta import relativedelta
from pptx import Presentation


def get_secrets():
    try:
        CLIENT_ID = os.getenv("CLIENT_ID")
        CLIENT_SECRET = os.getenv("CLIENT_SECRET")
        TENANT_ID = os.getenv("TENANT_ID")
        HUBSPOT_API_KEY = os.getenv("HUBSPOT_API_KEY")
        OPEN_AI_KEY = os.getenv("OPEN_AI_KEY")
        USER_ID = os.getenv("USER_ID")
        print(USER_ID)
        return (
            CLIENT_ID,
            CLIENT_SECRET,
            TENANT_ID,
            HUBSPOT_API_KEY,
            OPEN_AI_KEY,
            USER_ID,
        )
    except Exception as e:
        logging.error(f"Failed to retrieve secrets: {e}")
        raise

# from dotenv import load_dotenv
# mode = "dev"
# if mode == "dev":
#     load_dotenv()

(
    CLIENT_ID,
    CLIENT_SECRET,
    TENANT_ID,
    HUBSPOT_API_KEY,
    OPEN_AI_KEY,
    USER_ID,
) = get_secrets()
print("Env Setup")


def search_hubspot_object(object_type, search_body):
    url = f"https://api.hubapi.com/crm/v3/objects/{object_type}/search"
    headers = {
        "Content-Type": "application/json",
        "Authorization": f"Bearer {HUBSPOT_API_KEY}",
    }
    response = requests.post(url, headers=headers, data=json.dumps(search_body))
    response.raise_for_status()
    return response.json()


# Cache for storing owner details
owner_details_cache = {}


# Function to fetch owner details with throttling and caching
def fetch_owner_details(owner_id):
    if owner_id in owner_details_cache:
        return owner_details_cache[owner_id]
    url = f"https://api.hubapi.com/owners/v2/owners/{owner_id}"
    headers = {
        "Content-Type": "application/json",
        "Authorization": f"Bearer {HUBSPOT_API_KEY}",
    }
    retries = 3
    while retries > 0:
        try:
            response = requests.get(url, headers=headers)
            response.raise_for_status()
            owner_details_cache[owner_id] = response.json()
            return owner_details_cache[owner_id]
        except requests.exceptions.RequestException as e:
            logging.error(f"Failed to fetch owner details for ID {owner_id}: {e}")
            retries -= 1
            if retries == 0:
                raise
            time.sleep(2 ** (3 - retries))  # Exponential backoff


def fetch_and_attach_owner_details(deals, owner_property):
    owner_ids = [
        deal["properties"].get(owner_property)
        for deal in deals
        if deal["properties"].get(owner_property)
    ]
    unique_owner_ids = list(set(owner_ids))
    owners = {}
    for owner_id in unique_owner_ids:
        owners[owner_id] = fetch_owner_details(owner_id)
        time.sleep(0.1)  # Throttle requests to avoid hitting rate limits
    for deal in deals:
        owner_id = deal["properties"].get(owner_property)
        if owner_id:
            deal[f"{owner_property}_details"] = owners.get(owner_id, {})
    return deals


def fetch_notes(deal_id):
    url = f"https://api.hubapi.com/crm/v3/objects/notes/search"
    headers = {
        "Content-Type": "application/json",
        "Authorization": f"Bearer {HUBSPOT_API_KEY}",
    }
    search_body = {
        "filterGroups": [
            {
                "filters": [
                    {
                        "propertyName": "associations.deal",
                        "operator": "EQ",
                        "value": deal_id,
                    }
                ]
            }
        ],
        "properties": ["hs_note_body", "hs_attachment_ids"]
    }
    retries = 3
    while retries > 0:
        try:
            response = requests.post(url, headers=headers, data=json.dumps(search_body))
            response.raise_for_status()
            return response.json().get('results', [])
        except requests.exceptions.RequestException as e:
            retries -= 1
            if retries == 0:
                raise
            time.sleep(2 ** (3 - retries))  # Exponential backoff


def download_file(url):
    response = requests.get(url)
    response.raise_for_status()
    return response.content


def read_pdf(file_content):
    pdf_reader = PyPDF2.PdfReader(BytesIO(file_content))
    text = ""
    for page_num in range(len(pdf_reader.pages)):
        page = pdf_reader.pages[page_num]
        text += page.extract_text()
    return text


def read_word(file_content):
    doc = docx.Document(BytesIO(file_content))
    text = "\n".join([paragraph.text for paragraph in doc.paragraphs])
    return text


def read_excel(file_content):
    excel_file = pd.ExcelFile(BytesIO(file_content))
    sheets = {}
    for sheet_name in excel_file.sheet_names:
        sheets[sheet_name] = excel_file.parse(sheet_name)
    return sheets


def read_ppt(file_content):
    presentation = Presentation(BytesIO(file_content))
    text = ""
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text


def update_file_access(file_id, access_level):
    url = f"https://api.hubapi.com/files/v3/files/{file_id}"
    headers = {
        "Authorization": f"Bearer {HUBSPOT_API_KEY}",
        "Content-Type": "application/json"
    }
    data = {
        "options": {
            "access": access_level
        }
    }
    response = requests.patch(url, headers=headers, data=json.dumps(data))
    response.raise_for_status()
    return response.json()


def generate_signed_url(file_id):
    url = f"https://api.hubapi.com/files/v3/files/{file_id}/signed-url"
    headers = {
        "Authorization": f"Bearer {HUBSPOT_API_KEY}",
        "Content-Type": "application/json"
    }
    response = requests.get(url, headers=headers)
    response.raise_for_status()
    return response.json()["url"]


def get_file_details(file_id):
    url = f"https://api.hubapi.com/files/v3/files/{file_id}"
    headers = {
        "Authorization": f"Bearer {HUBSPOT_API_KEY}"
    }
    response = requests.get(url, headers=headers)
    response.raise_for_status()
    return response.json()


# def fetch_attachment(file_id):
#     url = f"https://api.hubapi.com/filemanager/api/v3/files/{file_id}/signed-url"
#     headers = {
#         "Content-Type": "application/json",
#         "Authorization": f"Bearer {HUBSPOT_API_KEY}",
#     }
#     retries = 3
#     while retries > 0:
#         try:
#             response = requests.get(url, headers=headers)
#             response.raise_for_status()
#             file = response.json()
#             file_url = file['url']
#             print("FILE--", file)
#             if file_url:
#                 file_content = download_file(file_url)
#                 return file_content, file_info["name"], signed_url
#             else:
#                 raise Exception(f"No URL found for file ID {file_id}")
#         except requests.exceptions.RequestException as e:
#             logging.error(f"Failed to fetch attachment for file ID {file_id}: {e}")
#             retries -= 1
#             if retries == 0:
#                 raise
#             time.sleep(2 ** (3 - retries))  # Exponential backoff


def attach_notes(deals):
    for deal in deals:
        # time.sleep(0.5)
        deal_id = deal['id']
        notes = fetch_notes(deal_id)
        deal['notes'] = notes
    return deals


failed_attachments = []


def attach_attachments(deals):
    for deal in deals:
        if deal['notes']:
            for note in deal['notes']:
                if note['properties']['hs_attachment_ids']:
                    file_id = note['properties']['hs_attachment_ids']
                    # time.sleep(0.4)
                    print("file_id", file_id)
                    try:
                        attachment = fetch_attachment(file_id)
                        deal['attachments'] = attachment
                    except:
                        failed_attachments.append(file_id)
                        continue
    return deals


def fetch_engagements(deal_id, engagement_type="EMAIL"):
    url = f"https://api.hubapi.com/engagements/v1/engagements/associated/deal/{deal_id}/paged"
    headers = {
        "Authorization": f"Bearer {HUBSPOT_API_KEY}",
        "Content-Type": "application/json"
    }
    params = {
        "limit": 100,
        "offset": 0
    }
    engagements = []
    while True:
        response = requests.get(url, headers=headers, params=params)
        response.raise_for_status()
        data = response.json()
        engagements.extend(
            [engagement for engagement in data['results'] if engagement['engagement']['type'] == engagement_type])
        if not data['hasMore']:
            break
        params['offset'] = data['offset']
    return engagements


def attach_engagements(deals):
    for deal in deals:
        time.sleep(0.5)
        deal_id = deal['id']
        engagements = fetch_engagements(deal_id)
        deal['engagements'] = engagements
    return deals


def fetch_deals(start_date=None, end_date=None):
    properties = [
        "dealname",
        "priority",
        "referral_type",
        "pipeline",
        "broad_category_updated",
        "subcategory",
        "fund",
        "hubspot_owner_id",
        "team_member_1",
        "createdate",
        "keywords",
    ]
    deals = []
    if start_date is None and end_date is None:
        start_date = str(datetime.now().date() + relativedelta(months=-3))
        end_date = str(datetime.now().date() + relativedelta(days=+1))
    start_date = datetime.strptime(start_date, "%Y-%m-%d")
    end_date = datetime.strptime(end_date, "%Y-%m-%d")
    # Hubspot needs datetime to be set to midnight
    start_date = start_date.replace(hour=0, minute=0, second=0, microsecond=0)
    end_date = end_date.replace(hour=0, minute=0, second=0, microsecond=0)
    # Convert the close date to a Unix timestamp in milliseconds
    start_date = int(start_date.timestamp() * 1000)
    end_date = int(end_date.timestamp() * 1000)
    after = None
    while True:
        search_body = {
            "filterGroups": [
                {
                    "filters": [
                        {
                            "propertyName": "createdate",
                            "operator": "BETWEEN",
                            "highValue": end_date,
                            "value": start_date,
                        },
                        {
                            "propertyName": "pipeline",
                            "operator": "EQ",
                            "value": "default",
                        }
                    ]
                }
            ],
            "properties": properties,
            "limit": 100,
        }
        if after:
            search_body["after"] = after
        search_results = search_hubspot_object("deals", search_body)
        deals.extend(search_results.get("results", []))
        pagination = search_results.get("paging", [])
        if pagination and "next" in pagination:
            after = pagination["next"]["after"]
        else:
            break
    return deals


def get_engagements_for_deal(deal_id):
    url = f"https://api.hubapi.com/engagements/v1/engagements/associated/deal/{deal_id}/paged"
    headers = {
        "Authorization": f"Bearer {HUBSPOT_API_KEY}",
        "Content-Type": "application/json"
    }
    response = requests.get(url, headers=headers)
    response.raise_for_status()
    return response.json()


def fetch_attachment(file_id):
    url = f"https://api.hubapi.com/filemanager/api/v3/files/{file_id}/signed-url"
    headers = {
        "Authorization": f"Bearer {HUBSPOT_API_KEY}",
        "Content-Type": "application/json",
        "User-Agent": "Custom"
    }
    response = requests.get(url, headers=headers)
    response.raise_for_status()
    file_url = response.json()["url"]
    extension = response.json()["extension"]
    name = response.json()["name"]
    file_response = requests.get(file_url)
    file_content = file_response.content
    if extension == "pdf":
        return read_pdf(file_content)
    elif extension == "docx":
        return read_word(file_content)
    elif extension == "xlsx":
        return read_excel(file_content)
    elif extension == "pptx":
        return read_ppt(file_content)
    else:
        raise ValueError(f"Unsupported file type: {extension} for {name}")


def organize_deals(deals):
    deals_dict = {}
    for deal in deals:
        broad_category = deal['properties']['broad_category_updated']
        if broad_category not in deals_dict:
            deals_dict[broad_category] = deal
    return deals_dict


def organize_cleaned_deals(cleaned_deals):
    deals_dict = {}
    for deal in cleaned_deals:
        broad_category = deal['broad_category_updated']
        if broad_category not in deals_dict:
            deals_dict[broad_category] = [deal]
        else:
            deals_dict[broad_category].append(deal)
    return deals_dict


def read_prompt_text(text_path):
    try:
        with open(text_path, "r") as file:
            return file.read()
    except FileNotFoundError:
        return "Flag as no prompt"


gpt_prompt_path = "data/gpt_prompt.txt"
gpt_prompt = read_prompt_text(gpt_prompt_path)


def create_openai_client(OPEN_AI_KEY):
    openai_client = openai.OpenAI(api_key=OPEN_AI_KEY)
    return openai_client


openai_client = create_openai_client(OPEN_AI_KEY)

gpt_errors = []


def parse_with_chatgpt(openai_client, deal):
    try:
        messages = [
            {"role": "system", "content": gpt_prompt},
            {
                "role": "user",
                "content": f"Deal info: {deal}",
            },
        ]
        response = openai_client.chat.completions.create(
            model="gpt-4o",
            messages=messages,
            max_tokens=2500,
            n=1,
            stop=None,
            temperature=0.5,
        )
        return response.choices[0].message.content
    except Exception as e:
        gpt_errors.append({deal['id']: str(e)})
        return None


user_prompt_for_final_path = "data/user_prompt_for_final.txt"
user_prompt_for_final = read_prompt_text(user_prompt_for_final_path)

system_prompt_for_final_path = "data/system_prompt_for_final.txt"
system_prompt_for_final = read_prompt_text(system_prompt_for_final_path)


# {"custom_id": "request-1", "method": "POST", "url": "/v1/chat/completions", "body": {"model": "gpt-3.5-turbo-0125", "messages": [{"role": "system", "content": "You are a helpful assistant."},{"role": "user", "content": "Hello world!"}],"max_tokens": 1000}}
# {"custom_id": "request-2", "method": "POST", "url": "/v1/chat/completions", "body": {"model": "gpt-3.5-turbo-0125", "messages": [{"role": "system", "content": "You are an unhelpful assistant."},{"role": "user", "content": "Hello world!"}],"max_tokens": 1000}}

def batch_with_chatgpt(openai_client, deals):
    jsonl_lines = []
    for deal in deals:
        prompt = {
                "custom_id": deal['properties']['dealname'],
                "method": "POST",
                "url": "/v1/chat/completions",
                "body": {
                    "model": "gpt-4o",
                    "messages": [
                        {"role": "system", "content": gpt_prompt},
                        {"role": "user", "content": f"Deal info: {deal}"}
                    ],
                    "max_tokens": 2500,
                    "n": 1,
                    "stop": None,
                    "temperature": 0.5
                    }
                }
        jsonl_lines.append(json.dumps(prompt))
    with open('deals_prompts.jsonl', 'w') as f:
        for line in jsonl_lines:
            f.write(line + '\n')
    batch_input_file = openai_client.files.create(
        file=open("deals_prompts.jsonl", "rb"),
        purpose="batch"
    )
    batch_input_file_id = batch_input_file.id
    batch = openai_client.batches.create(
        input_file_id=batch_input_file_id,
        endpoint="/v1/chat/completions",
        completion_window="24h",
        metadata={
            "description": "deal data recommendation generator"
        }
    )
    return batch


def check_gpt(openai_client, batch):
    # check status
    file_response = None
    retrieved_batch = openai_client.batches.retrieve(batch.id)
    if retrieved_batch.status == "completed":
        file_response = openai_client.files.content(retrieved_batch.output_file_id)
        return file_response.text

def delete_batch_file(openai_client, batch):
    # finish delete file

    all_batches = openai_client.batches.list(limit=10)
    # delete function
    files = openai_client.File.list()

    openai_client.files.delete(batch.output_file_id)


def compile_with_chatgpt(openai_client, cleaned_deals):
    try:
        messages = [
            {"role": "system", "content": system_prompt_for_final},
            {
                "role": "user",
                "content": f"Context: {user_prompt_for_final} - Input Data: {cleaned_deals}",
            },
        ]
        response = openai_client.chat.completions.create(
            model="gpt-4o",
            messages=messages,
            max_tokens=2500,
            n=1,
            stop=None,
            temperature=0.5,
        )
        return response.choices[0].message.content
    except Exception as e:
        print(f"Error: {e}")
        return None



def add_field_to_deal(field):
    url = "https://api.hubapi.com/properties/v1/deals/properties"

    headers = {
        'Content-Type': 'application/json',
        'Authorization': f'Bearer {HUBSPOT_API_KEY}',
    }

    label = field
    name = field.replace(' ', '_').lower()

    payload = {
        "name": "broad_category",
        "label": "Broad Category",
        "description": "The broad category of the deal",
        "groupName": "dealinformation",
        "type": "string",
        "fieldType": "text"
    }

    response = requests.post(url, json=payload, headers=headers)

    if response.status_code == 201:
        print("Custom property created successfully!")
    else:
        print(f"Failed to create property: {response.status_code}, {response.text}")

def update_deal(deal):
    update_url = f'https://api.hubapi.com/deals/v1/deal/{deal.id}'
    headers = {
        "Authorization": f"Bearer {HUBSPOT_API_KEY}",
        "Content-Type": "application/json"
    }

    update_payload = {
        "properties": [
            {
                "name": "broad_category_updated",
                "value": "value"
            },
            {
                "name": "subcategory",
                "value": "value"
            }
        ]
    }

    response = requests.put(update_url, json=update_payload, headers=headers)

    if response.status_code == 200:
        print("Deal updated successfully!")
    else:
        print(f"Failed to update deal: {response.status_code}, {response.text}")

def create_hubspot_field():
    pass

def update_hubspot_keywords(deal):
    update_url = f"https://api.hubapi.com/deals/v1/deal/{deal['id']}"
    headers = {
        "Authorization": f"Bearer {HUBSPOT_API_KEY}",
        "Content-Type": "application/json"
    }
    if deal.get('parsed'):
        if deal['parsed'].get('recommendation'):
            rec_keys = str(deal['parsed']['recommendation'])
            update_payload = {
                "properties": [
                    {
                        "name": "keywords",
                        "value": rec_keys
                    }
                ]
            }
            response = requests.put(update_url, json=update_payload, headers=headers)
            return response


def export_csv(start_date=None, end_date=None):
    deals = fetch_deals(start_date=start_date, end_date=end_date)
    deals = fetch_and_attach_owner_details(deals, "hubspot_owner_id")
    deals = fetch_and_attach_owner_details(deals, "team_member_1")
    flattened_data = []
    for deal in deals:
        flattened_entry = deal['properties'].copy()
        flattened_entry['id'] = deal['id']
        flattened_entry['createdAt'] = deal['createdAt']
        flattened_entry['updatedAt'] = deal['updatedAt']
        flattened_entry['archived'] = deal['archived']
        # Concatenate firstName and lastName for Lead Owner
        if deal.get('hubspot_owner_id_details'):
            lead_owner = deal.pop('hubspot_owner_id_details')
            flattened_entry['Lead Owner Name'] = f"{lead_owner['firstName']} {lead_owner['lastName']}"
            flattened_entry['Lead Owner Email'] = lead_owner['email']
        # Concatenate firstName and lastName for Support Member
        if deal.get('team_member_1_details'):
            support_member = deal.pop('team_member_1_details')
            flattened_entry['Support Member Name'] = f"{support_member['firstName']} {support_member['lastName']}"
            flattened_entry['Support Member Email'] = support_member['email']
        flattened_data.append(flattened_entry)

    # Convert to DataFrame
    df = pd.DataFrame(flattened_data)

    # Select only the necessary columns
    columns_to_include = ['broad_category_updated', 'createdate', 'dealname', 'fund', 'keywords',
           'pipeline', 'priority', 'referral_type', 'subcategory',
           'id', 'createdAt', 'updatedAt', 'archived', 'Lead Owner Name',
           'Lead Owner Email', 'Support Member Name', 'Support Member Email']
    #
    df = df[columns_to_include]

    if start_date is None and end_date is None:
        start_date = str(datetime.now().date() + relativedelta(months=-3))
        end_date = str(datetime.now().date() + relativedelta(days=+1))
    start_date = str(datetime.strptime(start_date, "%Y-%m-%d").date())
    end_date = str(datetime.strptime(end_date, "%Y-%m-%d").date())
    # Save to CSV
    file_object = io.StringIO()
    df.to_csv(file_object, index=False)
    file_object.seek(0)
    filename = f'Deal_Export--{start_date}-{end_date}.csv'

    return file_object, filename